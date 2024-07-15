import os
import argparse
import requests
from pptx import Presentation
from tqdm import tqdm

'''
You need to pip install to your python enviroment the following command:

pip install requests python-pptx tqdm
'''

# Set your Google Cloud API key
API_KEY = ""

# GET languages supported
def get_supported_languages(api_key):
    url = f"https://translation.googleapis.com/language/translate/v2/languages?key={api_key}"
    response = requests.get(url)
    if response.status_code == 200:
        languages = response.json().get("data", {}).get("languages", [])
        return [lang['language'] for lang in languages]
    else:
        print(f"Error fetching supported languages: {response.status_code} {response.text}")
        return []

# POST translate test
def translate_text(text, target_language):
    url = f"https://translation.googleapis.com/language/translate/v2?key={API_KEY}"
    headers = {
        "Content-Type": "application/json"
    }
    body = {
        "q": text,
        "target": target_language,
        "format": "text"
    }
    response = requests.post(url, headers=headers, json=body)
    if response.status_code == 200:
        return response.json()['data']['translations'][0]['translatedText']
    else:
        print(f"Error translating text: {response.status_code} {response.text}")
        return text

def translate_shape_text(shape, target_language):
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            translated_text = translate_text(run.text, target_language)
            run.text = translated_text

def process_presentation(input_file, target_language):
    print(f"Opening {input_file}")
    try:
        input_ppt = Presentation(input_file)
    except Exception as e:
        print(f"Error opening file {input_file}: {e}")
        return

    slide_count = len(input_ppt.slides)
    
    with tqdm(total=slide_count, desc="Translating", unit="slide") as pbar:
        for i, slide in enumerate(input_ppt.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    try:
                        translate_shape_text(shape, target_language)
                    except Exception as e:
                        print(f"Error processing shape on slide {i}: {e}")
            pbar.update(1)

    output_file = f"{target_language}_{input_file}"
    try:
        input_ppt.save(output_file)
        print(f"\nSaved as {output_file}")
    except Exception as e:
        print(f"Error saving file {output_file}: {e}")

def main():
    # for link output to console (this sometimes works)
    def link(uri, label=None):
        if label is None: 
            label = uri
        parameters = ''

        # OSC 8 ; params ; URI ST <name> OSC 8 ;; ST 
        escape_mask = '\033]8;{};{}\033\\{}\033]8;;\033\\'

        return escape_mask.format(parameters, uri, label)
    # ------------------------

    print("Please use the ISO 639 language code for the argument!")
    print("Example language syntax:")
    print("English: en")
    print("Spanish: es")
    print("French: fr")
    print("German: de")
    print("")
    print("See Full List of ISO 639 Languages here: " + link('https://cloud.google.com/translate/docs/languages'))
    print("")

    parser = argparse.ArgumentParser(description="Translate a PowerPoint presentation. Usage: python3 translatePPTX.py <input_pptx_file> <target_language>")
    parser.add_argument("input_file", nargs='?', help="Path to the input PowerPoint file")
    parser.add_argument("target_language", nargs='?', help="Target language for translation (ex: 'en' for English, 'es' for Spanish)")
    parser.add_argument("--list-langs", "-l", action="store_true", help="List supported languages and exit")
    args = parser.parse_args()

    # If --list-langs flag is provided, list the supported languages and exit
    if args.list_langs:
        supported_languages = get_supported_languages(API_KEY)
        if supported_languages:
            print("Supported languages:")
            print(supported_languages)
        else:
            print("Failed to fetch supported languages.")
        return

    if not args.input_file or not args.target_language:
        parser.print_help()
        return

    # check if users code is valid before running it (otherwise it will cause a bunch of errors @ api endpoint)
    valid_language_codes = get_supported_languages(API_KEY)

    if args.target_language in valid_language_codes:
        process_presentation(args.input_file, args.target_language)
    else:
        print("ERROR: NOT A VALID LANGUAGE CODE")
        print("")
        print("ERROR: Please submit a valid language code")
        print("")

        return -1


if __name__ == "__main__":
    main()
